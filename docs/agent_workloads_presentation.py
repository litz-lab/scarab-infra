#!/usr/bin/env python3
"""Generate PPTX presentation: CPU Microarchitecture Evaluation of Agentic AI Workloads."""

import json
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Helpers ──────────────────────────────────────────────────────────────
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
DARK  = RGBColor(0x2D, 0x2D, 0x2D)
BLUE  = RGBColor(0x1F, 0x77, 0xB4)
ORANGE = RGBColor(0xFF, 0x7F, 0x0E)
GREEN  = RGBColor(0x2C, 0xA0, 0x2C)
RED    = RGBColor(0xD6, 0x27, 0x28)
PURPLE = RGBColor(0x94, 0x67, 0xBD)
GRAY   = RGBColor(0x7F, 0x7F, 0x7F)
LIGHT_GRAY = RGBColor(0xF0, 0xF0, 0xF0)
HEADER_BG = RGBColor(0x2B, 0x57, 0x9A)

def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_title_shape(slide, text, left, top, width, height, font_size=28, bold=True, color=DARK):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    return tf

def add_body_text(slide, text, left, top, width, height, font_size=14, color=DARK):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(text.split('\n')):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.space_after = Pt(4)
    return tf

def add_bullet_list(slide, items, left, top, width, height, font_size=14, color=DARK):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        if isinstance(item, tuple):
            # (text, indent_level)
            p.text = item[0]
            p.level = item[1]
        else:
            p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.space_after = Pt(3)
    return tf

def _set_cell_borders(cell, color="000000", width="12700"):
    """Set all four borders of a table cell."""
    from lxml import etree
    from pptx.oxml.ns import qn
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    for border_name in ("lnL", "lnR", "lnT", "lnB"):
        # Remove existing border element
        for existing in tcPr.findall(qn(f"a:{border_name}")):
            tcPr.remove(existing)
        # Build border XML directly
        border_xml = (
            f'<a:{border_name} xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"'
            f' w="{width}" cap="flat" cmpd="sng" algn="ctr">'
            f'<a:solidFill><a:srgbClr val="{color}"/></a:solidFill>'
            f'<a:prstDash val="solid"/>'
            f'</a:{border_name}>'
        )
        border_elem = etree.fromstring(border_xml)
        tcPr.append(border_elem)

def add_table(slide, data, left, top, width, height, header_color=None):
    rows, cols = len(data), len(data[0])
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    # Remove default blue table style
    from lxml import etree
    tbl = table._tbl
    tblPr = tbl.find('{http://schemas.openxmlformats.org/drawingml/2006/main}tblPr')
    if tblPr is not None:
        tblPr.attrib.clear()

    for col_idx in range(cols):
        for row_idx in range(rows):
            cell = table.cell(row_idx, col_idx)
            cell.text = str(data[row_idx][col_idx])
            cell.fill.background()
            _set_cell_borders(cell, color="000000", width="12700")
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(10)
                paragraph.font.color.rgb = DARK
                paragraph.alignment = PP_ALIGN.CENTER
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            if row_idx == 0:
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.bold = True
    return table

# ── Load perf data ──────────────────────────────────────────────────────
with open('/soe/surim/src/infra_agent_perf/workloads/workloads_db.json') as f:
    db = json.load(f)
agent_data = db.get('agent', {})

# ── Build Presentation ──────────────────────────────────────────────────
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

L_MARGIN = Inches(0.7)
CONTENT_W = Inches(12)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 1: Title
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(slide, WHITE)

add_title_shape(slide, "CPU Microarchitecture Evaluation\nof Agentic AI Workloads",
                L_MARGIN, Inches(1.5), CONTENT_W, Inches(2.5), font_size=36, color=HEADER_BG)
add_body_text(slide,
    "SimPoint-based trace collection & Scarab simulation\n"
    "for LangChain, LangGraph, AutoGen, Haystack (RAG), and SWE-Agent",
    L_MARGIN, Inches(3.8), CONTENT_W, Inches(1.5), font_size=20, color=GRAY)
add_body_text(slide, "Surim Oh  |  April 2026", L_MARGIN, Inches(5.5), CONTENT_W, Inches(0.5),
              font_size=16, color=GRAY)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 2: Motivation & Overview
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)
add_title_shape(slide, "Motivation", L_MARGIN, Inches(0.3), CONTENT_W, Inches(0.8), font_size=30, color=HEADER_BG)

add_bullet_list(slide, [
    "Agentic AI is an emerging CPU workload class — orchestration, tool calling, state management",
    "Current CPU benchmarks (SPEC, Datacenter) don't capture agent execution patterns",
    "Agent workloads have distinct microarchitectural signatures:",
    ("Frontend-bound: deep call stacks, polymorphic dispatch, large code footprints", 1),
    ("Backend-bound: JSON/string processing, hash table lookups, dict manipulation", 1),
    ("Branch-heavy: conditional tool selection, state-machine transitions", 1),
    "Goal: characterize 5 agent workloads spanning frameworks and application archetypes",
    "Method: perf stat (real HW) + SimPoint traces → Scarab simulation",
], L_MARGIN, Inches(1.2), CONTENT_W, Inches(5.5), font_size=16)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 3: Workload Overview Table
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)
add_title_shape(slide, "Workload Suite Overview", L_MARGIN, Inches(0.3), CONTENT_W, Inches(0.8),
                font_size=30, color=HEADER_BG)

table_data = [
    ["Workload", "Framework", "Archetype", "Key Computation", "Mock Components", "Peak RSS"],
    ["LangChain", "LangChain Core", "Simple agent loop", "JSON ser/deser, string ops", "Fake tool (local_search)", "50 MB"],
    ["LangGraph", "LangGraph", "ReAct agent (graph)", "JSON, message passing, graph traversal", "FakeToolCallingModel, 3 tools", "~80 MB"],
    ["AutoGen", "AutoGen AgentChat", "Multi-agent (async)", "Async orchestration, JSON, dict ops", "ReplayChatCompletionClient", "~90 MB"],
    ["Haystack", "Sentence-Transformers\n+ FAISS", "RAG pipeline", "FP32 embedding, SIMD dot-product,\nnumpy matrix ops", "Real model, synthetic corpus", "519 MB"],
    ["SWE-Agent", "LangGraph StateGraph", "Multi-step planning\nagent", "Regex, string/dict manipulation,\nbranchy control flow", "6 mock dev tools, scripted\n20-step sequence", "252 MB"],
]
add_table(slide, table_data, L_MARGIN, Inches(1.3), Inches(12), Inches(4.5))

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 4: Workload Algorithms - LangChain & LangGraph
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)
add_title_shape(slide, "Workload Algorithms (1/3): LangChain & LangGraph",
                L_MARGIN, Inches(0.3), CONTENT_W, Inches(0.8), font_size=28, color=HEADER_BG)

add_body_text(slide, "LangChain — Simple Agent Loop", L_MARGIN, Inches(1.2), Inches(5.5), Inches(0.5),
              font_size=18, color=BLUE)
add_bullet_list(slide, [
    "Pipeline: plan → select_tool → marshal_request → invoke_tool → parse_response → synthesize",
    "Single tool (local_search) returns deterministic JSON results",
    "Context blob scales with --context-size: 416B (short) / 1.6KB (med) / 6.6KB (long)",
    "Computational core: JSON encode/decode with sort_keys, dict operations",
], L_MARGIN, Inches(1.7), Inches(5.5), Inches(2.5), font_size=13)

add_body_text(slide, "LangGraph — Graph-Structured ReAct Agent", Inches(7), Inches(1.2), Inches(5.5), Inches(0.5),
              font_size=18, color=BLUE)
add_bullet_list(slide, [
    "Uses create_react_agent with FakeToolCallingModel",
    "Per iteration: 3 tool calls (cpu_analyze → memory_profile → branch_analyze)",
    "Adds message accumulation and graph-node traversal overhead",
    "More framework overhead than LangChain: state checkpointing, edge routing",
], Inches(7), Inches(1.7), Inches(5.5), Inches(2.5), font_size=13)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 5: Workload Algorithms - AutoGen
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)
add_title_shape(slide, "Workload Algorithms (2/3): AutoGen & Haystack",
                L_MARGIN, Inches(0.3), CONTENT_W, Inches(0.8), font_size=28, color=HEADER_BG)

add_body_text(slide, "AutoGen — Async Multi-Agent Framework", L_MARGIN, Inches(1.2), Inches(5.5), Inches(0.5),
              font_size=18, color=GREEN)
add_bullet_list(slide, [
    "AssistantAgent with ReplayChatCompletionClient (7 scripted responses)",
    "3 async tools: topdown_analyze, cache_profile, branch_stats",
    "Python asyncio event loop adds coroutine scheduling overhead",
    "Distinct from LangChain/LangGraph: async/await dispatch, different runtime path",
], L_MARGIN, Inches(1.7), Inches(5.5), Inches(2.5), font_size=13)

add_body_text(slide, "Haystack — RAG Pipeline (Real ML Inference)", Inches(7), Inches(1.2), Inches(5.5), Inches(0.5),
              font_size=18, color=GREEN)
add_bullet_list(slide, [
    "Real SentenceTransformer model: all-MiniLM-L6-v2 (22M params, FP32)",
    "Pipeline: embed_query → FAISS vector_search (top-k=10) → rerank → synthesize",
    "Corpus scales: 100 docs (short) / 500 (med) / 2000 (long)",
    "SIMD-heavy: PyTorch FP32 inference, FAISS IndexFlatIP dot products, numpy ops",
    "Only workload with real ML compute — expected backend-bound (memory/FP)",
], Inches(7), Inches(1.7), Inches(5.5), Inches(3), font_size=13)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 6: Workload Algorithms - SWE-Agent
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)
add_title_shape(slide, "Workload Algorithms (3/3): SWE-Agent",
                L_MARGIN, Inches(0.3), CONTENT_W, Inches(0.8), font_size=28, color=HEADER_BG)

add_body_text(slide, "SWE-Agent — Multi-Step Planning Agent (Bug-Fix Simulation)",
              L_MARGIN, Inches(1.2), CONTENT_W, Inches(0.5), font_size=18, color=PURPLE)
add_bullet_list(slide, [
    "Built on LangGraph StateGraph with 5 nodes: planner → tool_selector → tool_executor → observer → check_done",
    "Rich AgentState: task, plan[], completed[], file_state{}, observations[], step_count, messages",
    "6 mock developer tools (pure-CPU, in-memory file system):",
    ("read_file, write_file, list_dir — dict key/value operations", 1),
    ("grep_files — regex compilation + search across file corpus", 1),
    ("run_tests — conditional pass/fail based on step count", 1),
    ("apply_patch — deterministic string mutation (simulates a-b → a+b fix)", 1),
    "20-step scripted sequence per iteration with conditional replanning every 2 steps",
    "Deepest control flow of all workloads — closest to real multi-step agent planning",
], L_MARGIN, Inches(1.8), CONTENT_W, Inches(4.5), font_size=14)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 7: Perf Characteristics — Top-Down Table
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)
add_title_shape(slide, "CPU Performance Characteristics (perf stat, Top-Down L1)",
                L_MARGIN, Inches(0.3), CONTENT_W, Inches(0.8), font_size=28, color=HEADER_BG)

# Build average topdown per framework
perf_table = [["Framework", "Retiring %", "Bad Spec %", "Frontend\nBound %", "Backend\nBound %", "Exec Time\n(representative)", "Peak RSS"]]
for fw_name, display, exec_time, rss in [
    ("langchain", "LangChain", "12s (short/250)", "50 MB"),
    ("langgraph", "LangGraph", "~2 min", "~80 MB"),
    ("autogen", "AutoGen", "~2 min", "~90 MB"),
    ("haystack", "Haystack", "2m31s (short/100)", "519 MB"),
    ("swe_agent", "SWE-Agent", "2m8s (short/150)", "252 MB"),
]:
    variants = agent_data.get(fw_name, {})
    rets, bads, fes, bes = [], [], [], []
    for v in variants.values():
        td = v.get("performance", {}).get("topdown", {})
        rets.append(td.get("retiring", 0))
        bads.append(td.get("bad_speculation", 0))
        fes.append(td.get("frontend_bound", 0))
        bes.append(td.get("backend_bound", 0))
    n = len(rets) or 1
    perf_table.append([
        display,
        f"{sum(rets)/n:.1f}",
        f"{sum(bads)/n:.1f}",
        f"{sum(fes)/n:.1f}",
        f"{sum(bes)/n:.1f}",
        exec_time,
        rss,
    ])

add_table(slide, perf_table, L_MARGIN, Inches(1.3), Inches(12), Inches(3.5))

add_bullet_list(slide, [
    "All agent workloads are dominated by backend-bound stalls (44-53%) — memory/data dependencies",
    "Frontend-bound is significant (25-32%) — large Python code footprint, indirect dispatch",
    "Retiring is low (17-22%) — typical of interpreted languages with high overhead",
    "Bad speculation is modest (4-5%) — branch prediction handles Python's patterns reasonably",
    "Haystack shows higher variance across variants — real ML inference vs framework overhead",
], L_MARGIN, Inches(5.0), CONTENT_W, Inches(2.2), font_size=14)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 8: Top-Down Stacked Bar Chart
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)
add_title_shape(slide, "Top-Down L1 Breakdown by Workload (Average Across Variants)",
                L_MARGIN, Inches(0.3), CONTENT_W, Inches(0.8), font_size=28, color=HEADER_BG)

# Build data for shape-based stacked bar chart (Google Slides compatible)
categories = []
retiring_vals, bad_vals, fe_vals, be_vals = [], [], [], []
for fw_name, display in [("langchain", "LangChain"), ("langgraph", "LangGraph"),
                          ("autogen", "AutoGen"), ("haystack", "Haystack"), ("swe_agent", "SWE-Agent")]:
    categories.append(display)
    variants = agent_data.get(fw_name, {})
    rets, bads, fes, bes = [], [], [], []
    for v in variants.values():
        td = v.get("performance", {}).get("topdown", {})
        rets.append(td.get("retiring", 0))
        bads.append(td.get("bad_speculation", 0))
        fes.append(td.get("frontend_bound", 0))
        bes.append(td.get("backend_bound", 0))
    n = len(rets) or 1
    retiring_vals.append(sum(rets)/n)
    bad_vals.append(sum(bads)/n)
    fe_vals.append(sum(fes)/n)
    be_vals.append(sum(bes)/n)

# Draw stacked bars using rectangle shapes (compatible with Google Slides)
bar_colors = [
    ("Retiring", RGBColor(0x2C,0xA0,0x2C)),
    ("Bad Speculation", RGBColor(0xFF,0x7F,0x0E)),
    ("Frontend Bound", RGBColor(0x1F,0x77,0xB4)),
    ("Backend Bound", RGBColor(0xD6,0x27,0x28)),
]
all_vals = list(zip(retiring_vals, bad_vals, fe_vals, be_vals))

chart_left = Inches(2.0)
chart_top = Inches(1.5)
chart_height = Inches(4.5)
bar_width = Inches(1.4)
bar_gap = Inches(0.5)
max_val = 100.0  # percentages sum to ~100

# Y-axis labels (0%, 25%, 50%, 75%, 100%)
for pct in [0, 25, 50, 75, 100]:
    y = chart_top + chart_height - (pct / max_val) * chart_height
    add_body_text(slide, f"{pct}%", Inches(1.2), y - Pt(7), Inches(0.7), Inches(0.3),
                  font_size=10, color=GRAY)
    # Gridline
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        chart_left, int(y), Inches(9.5), Pt(0.5))
    line.fill.solid()
    line.fill.fore_color.rgb = LIGHT_GRAY
    line.line.fill.background()

for i, (cat, vals) in enumerate(zip(categories, all_vals)):
    x = chart_left + i * (bar_width + bar_gap)
    y_bottom = chart_top + chart_height
    for j, (val, (name, color)) in enumerate(zip(vals, bar_colors)):
        seg_height = (val / max_val) * chart_height
        if seg_height < Pt(1):
            continue
        y_top = y_bottom - seg_height
        rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, int(y_top), bar_width, int(seg_height))
        rect.fill.solid()
        rect.fill.fore_color.rgb = color
        rect.line.color.rgb = WHITE
        rect.line.width = Pt(0.5)
        # Add percentage label inside segment if tall enough
        if seg_height > Pt(16):
            tf = rect.text_frame
            tf.word_wrap = False
            p = tf.paragraphs[0]
            p.text = f"{val:.0f}%"
            p.font.size = Pt(9)
            p.font.color.rgb = WHITE
            p.font.bold = True
            p.alignment = PP_ALIGN.CENTER
            tf.paragraphs[0].space_before = Pt(0)
        y_bottom = y_top
    # Category label below bar
    add_body_text(slide, cat, x, chart_top + chart_height + Inches(0.1),
                  bar_width, Inches(0.4), font_size=12, color=DARK)

# Legend
legend_top = Inches(1.3)
legend_left = Inches(10.0)
for j, (name, color) in enumerate(bar_colors):
    ly = legend_top + j * Inches(0.35)
    swatch = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, legend_left, ly, Inches(0.25), Inches(0.2))
    swatch.fill.solid()
    swatch.fill.fore_color.rgb = color
    swatch.line.fill.background()
    add_body_text(slide, name, legend_left + Inches(0.35), ly - Pt(2),
                  Inches(2), Inches(0.3), font_size=11, color=DARK)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 9: Framework Updates - scarab-infra
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)
add_title_shape(slide, "Infrastructure Updates: scarab-infra",
                L_MARGIN, Inches(0.3), CONTENT_W, Inches(0.8), font_size=28, color=HEADER_BG)

add_bullet_list(slide, [
    "Added 5 agent workloads with Docker-containerized environments:",
    ("LangChain, LangGraph, AutoGen — Python 3.10 + framework packages", 1),
    ("Haystack — sentence-transformers + FAISS + numpy (model pre-downloaded at build time)", 1),
    ("SWE-Agent — LangGraph StateGraph with rich state and 6 mock developer tools", 1),
    "Each workload: 6 variants (short/medium/long × 2 iteration counts) = 30 configurations",
    "Trace pipeline: cluster_then_trace (SimPoint) with libfpg.so fingerprinting",
    "Perf pipeline: perf stat with Top-Down L1 counters on pinned core",
    "Slurm integration: parallel job submission, auto image build on nodes",
    "Status tracking: ./sci --status shows per-job stage (fingerprinting/clustering/tracing/...)",
    "Key files modified:",
    ("workloads/agent/ — Dockerfile + run scripts (run_langchain.py, run_langgraph.py, etc.)", 1),
    ("json/agent_trace.json, agent_perf.json — 30-variant descriptors", 1),
    ("common/scripts/run_simpoint_trace.py — multi-thread handling, DR config isolation", 1),
    ("scripts/slurm_runner.py — stage progress, kill wait, image cleanup", 1),
], L_MARGIN, Inches(1.2), CONTENT_W, Inches(5.5), font_size=14)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 10: DynamoRIO Fixes
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)
add_title_shape(slide, "DynamoRIO Fixes (fork: 5surim/dynamorio)",
                L_MARGIN, Inches(0.3), CONTENT_W, Inches(0.8), font_size=28, color=HEADER_BG)

add_body_text(slide, "Two problems with DynamoRIO on multi-threaded Python workloads",
              L_MARGIN, Inches(1.2), CONTENT_W, Inches(0.5), font_size=16, color=RED)

add_bullet_list(slide, [
    "Root Cause 1: SIGSEGV crash — race in module_list_remove() (core/module_list.c)",
    ("module_list_remove() released write lock for client callback, module still in vmvector", 1),
    ("Another thread gets module name pointer via vmvector_lookup during this window", 1),
    ("First thread re-acquires lock, frees module data → second thread's pointer is stale → SIGSEGV", 1),
    "",
    "Root Cause 2: libfpg.so not loading — DR config cross-contamination on shared NFS",
    ("DR caches per-process configs in ~/.dynamorio/ keyed by binary+PID", 1),
    ("Concurrent slurm jobs on same NFS home pick up each other's config files", 1),
    ("Wrong config → libfpg.so client not loaded → empty fingerprint output", 1),
    "",
    "Fix 1 (for RC1): Root cause fix in module_list_remove() (core/module_list.c)",
    ("Move module_area_delete() + vmvector removal before releasing write lock", 1),
    ("Client callback only uses client_data (a copy), does not need original module_area_t", 1),
    ("Eliminates the race window entirely — no thread can get a stale pointer", 1),
    "",
    "Fix 2 (for RC1): strhash_key_cmp safe_read guard (core/hashtable.c)",
    ("Defense-in-depth: safe_read check on hash table key pointers before strcmp", 1),
    "",
    "Fix 3 (for RC2): Per-job HOME override (scarab-infra, run_simpoint_trace.py)",
    ("HOME={workload_dir} on drrun commands — isolates ~/.dynamorio/ per job", 1),
], L_MARGIN, Inches(1.7), CONTENT_W, Inches(5.2), font_size=13)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 11: Scarab Fix
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)
add_title_shape(slide, "Scarab Fix + Multi-Threaded Workload Handling",
                L_MARGIN, Inches(0.3), CONTENT_W, Inches(0.8), font_size=28, color=HEADER_BG)

add_body_text(slide, "Scarab Build Fix", L_MARGIN, Inches(1.2), Inches(5.5), Inches(0.5),
              font_size=18, color=BLUE)
add_bullet_list(slide, [
    "GCC 11 (Ubuntu 22.04) treats -Wclass-memaccess as error with -Werror",
    "cbp_tagescl_64k.h: memset() on non-trivial C++ classes (cbp64_folded_history, cbp64_lentry)",
    "Fix: replaced memset with default-construction loops:",
    ("for (auto &x : ch_i) x = cbp64_folded_history();", 1),
    ("for (auto &row : ch_t) for (auto &x : row) x = cbp64_folded_history();  // 2D array", 1),
    ("for (auto &x : ltable) x = cbp64_lentry();", 1),
], L_MARGIN, Inches(1.7), Inches(5.5), Inches(3), font_size=13)

add_body_text(slide, "Multi-Threaded Fingerprint Handling", Inches(7), Inches(1.2), Inches(5.5), Inches(0.5),
              font_size=18, color=BLUE)
add_bullet_list(slide, [
    "Python workloads spawn helper threads (GC, import, signal handlers)",
    "libfpg.so creates per-thread fingerprint files (bbfp.<thread_id>)",
    "cluster_then_trace expected exactly 1 bbfp file — failed with multiple",
    "Fix: pick the main thread's fingerprint (file with most segments)",
    "CPython GIL ensures main thread dominates instruction execution",
    "Simulating the dominant thread on a single core is representative",
], Inches(7), Inches(1.7), Inches(5.5), Inches(3), font_size=13)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 12: Current Status
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)
add_title_shape(slide, "Current Status", L_MARGIN, Inches(0.3), CONTENT_W, Inches(0.8),
                font_size=30, color=HEADER_BG)

status_table = [
    ["Component", "Status", "Details"],
    ["Perf collection\n(perf stat)", "Complete", "30/30 variants collected\nTop-Down L1 data in workloads_db.json"],
    ["Trace collection\n(SimPoint fingerprint)", "In progress", "30 jobs running on slurm cluster\n(cluster_then_trace mode with libfpg.so)"],
    ["Trace collection\n(SimPoint clustering + tracing)", "Pending", "Runs after fingerprinting completes"],
    ["Scarab simulation", "Pending", "Will run after traces are collected"],
    ["DynamoRIO fork", "Fixed", "Root cause fix in module_list_remove()\n+ strhash_key_cmp guard + HOME isolation"],
    ["Scarab fork", "Fixed", "GCC 11 build fix for cbp_tagescl_64k.h"],
]
add_table(slide, status_table, L_MARGIN, Inches(1.3), Inches(12), Inches(5))

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 13: Next Steps
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)
add_title_shape(slide, "Next Steps (1/2): Trace Collection & Analysis",
                L_MARGIN, Inches(0.3), CONTENT_W, Inches(0.8), font_size=30, color=HEADER_BG)

add_bullet_list(slide, [
    "1. Complete SimPoint trace collection (30 workload variants)",
    ("Fingerprinting → clustering → representative segment tracing", 1),
    ("Expected: ~10-20 simpoints per workload variant", 1),
    "",
    "2. Run Scarab simulation on collected traces",
    ("Simulate each simpoint segment with detailed microarchitecture model (add missing insts)", 1),
    ("Collect IPC, cache miss rates, branch misprediction rates, pipeline utilization", 1),
    "",
    "3. Detailed microarchitectural analysis",
    ("Compare agent workloads vs traditional benchmarks (SPEC, datacenter)", 1),
    ("Identify bottlenecks: i-cache misses, d-cache pressure, branch misprediction", 1),
    ("Characterize per-workload signatures: RAG (backend/SIMD) vs framework (frontend)", 1),
    "",
    "4. Control flow graph (CFG) analysis",
    ("Extract CFGs from traces: basic block size distribution, branch type breakdown", 1),
    ("Compare loop depth, fan-out, indirect branch density across workloads", 1),
    ("Hot code footprint: what fraction of unique PCs covers 90%/99% of execution?", 1),
], L_MARGIN, Inches(1.2), CONTENT_W, Inches(5.5), font_size=15)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 14: Next Steps (2/2) — Architecture Sensitivity Studies
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)
add_title_shape(slide, "Next Steps (2/2): Architecture Sensitivity Studies",
                L_MARGIN, Inches(0.3), CONTENT_W, Inches(0.8), font_size=30, color=HEADER_BG)

add_bullet_list(slide, [
    "5. Frontend-bound studies",
    ("I-cache size sweep (32/64/128KB) — Python's large code footprint likely causes high i-miss", 1),
    ("BTB size sensitivity — many indirect branch targets from polymorphic dispatch", 1),
    ("Fetch/decode width sensitivity — wider frontend vs capacity-bound?", 1),
    "",
    "6. Backend-bound studies",
    ("D-cache hierarchy sweep: L1D (32/48/64KB), L2 (256/512KB/1MB), LLC (2/4/8MB)", 1),
    ("Reuse distance analysis — quantify temporal locality in dict/JSON/hash table accesses", 1),
    ("ROB size sensitivity (256/384/512) — larger ROB to extract MLP and hide cache misses", 1),
    ("Prefetcher effectiveness — stride vs complex; irregular access patterns may defeat stride", 1),
    "",
    "7. Cross-cutting studies",
    ("Core width sweep (4/6/8-wide) — are agent workloads ILP-limited?", 1),
    ("Branch predictor config sweep — TAGE history length, indirect predictor sizing", 1),
    ("Haystack (FP/SIMD) vs framework (branchy/string) — quantify architectural divergence", 1),
    "",
    "8. Upstream DynamoRIO fix",
    ("Submit PR with module_list_remove() race fix to DynamoRIO upstream", 1),
], L_MARGIN, Inches(1.2), CONTENT_W, Inches(5.5), font_size=15)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 15: SimPoint Methodology
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)
add_title_shape(slide, "Appendix: SimPoint Methodology (cluster_then_trace)",
                L_MARGIN, Inches(0.3), CONTENT_W, Inches(0.8), font_size=28, color=HEADER_BG)

add_bullet_list(slide, [
    "Phase 1: Fingerprinting",
    ("Run workload under DynamoRIO with libfpg.so client", 1),
    ("Collect basic block frequency vectors every 10M instructions (segment)", 1),
    ("Output: bbfp.<tid> (fingerprint) + pcmap.<tid> (PC mapping)", 1),
    "",
    "Phase 2: Clustering",
    ("SimPoint k-means clustering on BB frequency vectors", 1),
    ("Select representative segments (simpoints) with weights", 1),
    ("Typically 10-20 clusters per workload", 1),
    "",
    "Phase 3: Trace Collection",
    ("Re-run workload under DynamoRIO drcachesim, trace only selected segments", 1),
    ("-trace_after_instrs N -trace_for_instrs M to capture specific instruction windows", 1),
    ("5 segments of warmup (50M instructions) before each simpoint", 1),
    "",
    "Phase 4: Simulation (next step)",
    ("Feed traces to Scarab cycle-accurate simulator", 1),
    ("Weighted IPC = sum(simpoint_weight × simpoint_IPC)", 1),
], L_MARGIN, Inches(1.2), CONTENT_W, Inches(5.5), font_size=14)

# ── Save ─────────────────────────────────────────────────────────────
output_path = "/soe/surim/src/infra_agent_perf/docs/agent_cpu_evaluation.pptx"
prs.save(output_path)
print(f"Saved to {output_path}")
